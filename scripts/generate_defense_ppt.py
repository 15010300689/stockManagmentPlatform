# -*- coding: utf-8 -*-
"""根据毕业设计论文要点生成约 5 分钟答辩用 PPT（需已安装 python-pptx）。"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _set_title(slide, text: str) -> None:
    tf = slide.shapes.title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)


def _bullet_body(slide, bullets: list[str]) -> None:
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, t in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(17)
        p.space_after = Pt(4)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1 封面
    slide_layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(slide_layout)
    box = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "库存管理系统设计与实现"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "毕业设计答辩（约 5 分钟）"
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "答辩人：贺燕珍\n北京理工大学继续教育学院 · 计算机科学与技术"
    p3.font.size = Pt(16)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)

    layout_title = prs.slide_layouts[1]

    slides_data = [
        (
            "一、研究背景与意义",
            [
                "数字化转型下，库存是供应链核心环节；传统台账/表格易出现账实不符、效率低、数据滞后。",
                "目标：设计并实现基于 Web 的库存管理系统，服务中小企业商品维护、出入库、分仓分位与预警统计。",
                "意义：提升数据准确性与周转效率，为管理决策提供数据支撑。",
            ],
        ),
        (
            "二、研究现状与本文定位",
            [
                "国外：智能化、云化、与 AI 结合较深，但对小规模场景存在成本高、过度设计问题。",
                "国内：Spring Boot 等轻量方案贴合中小企业；SaaS 渗透提升，仍需易落地、可扩展的自建系统。",
                "本文：采用 Spring Boot + MySQL + 前后端分离，聚焦轻量可维护的库存数字化方案。",
            ],
        ),
        (
            "三、需求与功能模块",
            [
                "功能：用户认证与权限（RBAC）、商品 CRUD 与检索、仓库/仓位分层、入库/出库/调整与流水、统计与低库存预警。",
                "非功能：界面可用性、JWT 鉴权、模块化可扩展、中小数据规模下的响应性能。",
            ],
        ),
        (
            "四、技术选型（与实现一致）",
            [
                "后端：Spring Boot 2.7、Spring MVC、MyBatis、MySQL、JWT 拦截与统一异常处理。",
                "前端：React 18、TypeScript、Vite、Ant Design 6；路由守卫与 Token 请求封装。",
                "工程：前后端分离，生产构建静态资源可由 Spring Boot 托管。",
            ],
        ),
        (
            "五、系统总体架构",
            [
                "表现层：React + Ant Design 页面（商品、仓库仓位、库存、统计预警、权限菜单等）。",
                "接口与安全层：Controller + JWT 校验，统一拦截非法请求。",
                "业务层：Service 组织商品、库存、菜单权限等逻辑，关键库存操作使用事务。",
                "持久层：MyBatis Mapper + XML；数据层 MySQL 存储业务与权限数据。",
            ],
        ),
        (
            "六、核心设计与实现要点",
            [
                "登录：校验用户名密码后签发 JWT；前端请求头携带 Token，拦截器统一校验。",
                "库存：入库增加库存并记流水；出库前校验数量，不足则拒绝；多表更新同一事务保证一致。",
                "仓位：parent_id 维护层级，后端组装树形结构供前端展示。",
            ],
        ),
        (
            "七、测试与验证",
            [
                "环境：Windows/macOS、JDK 8+、MySQL 8.x、Chrome、Postman 等。",
                "功能：登录与 Token、商品增删改查、入出库与库存联动、仓库仓位与预警列表等用例验证通过。",
                "结论：满足中小规模场景下的日常使用；高并发与压力测试可作为后续工作。",
            ],
        ),
        (
            "八、总结与展望",
            [
                "已完成：库存业务数字化、分仓分位、统计预警、基础 RBAC 与前后端分离架构。",
                "展望：更细粒度接口权限、采购销售与审批流、缓存与消息队列、可视化大屏与智能补货建议。",
            ],
        ),
    ]

    for title, bullets in slides_data:
        sl = prs.slides.add_slide(layout_title)
        _set_title(sl, title)
        _bullet_body(sl, bullets)

    # 结束页
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    box = sl.shapes.add_textbox(Inches(2), Inches(2.2), Inches(6), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢各位老师！"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "敬请批评指正"
    p2.font.size = Pt(22)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    out = (
        Path(__file__).resolve().parent.parent
        / "答辩PPT-库存管理系统设计与实现-约5分钟.pptx"
    )
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
